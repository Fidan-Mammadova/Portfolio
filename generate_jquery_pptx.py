from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "presentations" / "jquery_lessons_ru"
ASSETS_DIR = OUTPUT_DIR / "assets_v2"

THEME = {
    "bg": RGBColor(245, 248, 255),
    "title": RGBColor(31, 61, 122),
    "text": RGBColor(40, 45, 60),
    "accent": RGBColor(79, 113, 255),
    "muted": RGBColor(112, 120, 150),
}


def get_font(size: int):
    for name in ["Arial.ttf", "Helvetica.ttc", "DejaVuSans.ttf"]:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def draw_card(draw, x1, y1, x2, y2, title, lines):
    draw.rounded_rectangle([(x1, y1), (x2, y2)], radius=22, fill=(250, 252, 255), outline=(117, 140, 216), width=3)
    draw.rectangle([(x1, y1), (x2, y1 + 58)], fill=(79, 113, 255))
    draw.text((x1 + 18, y1 + 14), title, fill=(255, 255, 255), font=get_font(28))
    y = y1 + 78
    for line in lines:
        draw.text((x1 + 18, y), f"- {line}", fill=(37, 52, 87), font=get_font(24))
        y += 38


def create_infographic_image(path: Path, module_title: str, slide_title: str, visual_type: str):
    # RU: Рисуем обучающие схемы вместо случайных фото.
    # EN: Draw learning infographics instead of random photos.
    img = Image.new("RGB", (1200, 900), (232, 240, 255))
    draw = ImageDraw.Draw(img)

    draw.rectangle([(0, 0), (1200, 120)], fill=(36, 70, 170))
    draw.text((32, 35), module_title, fill=(255, 255, 255), font=get_font(34))
    draw.text((32, 140), slide_title, fill=(31, 61, 122), font=get_font(40))

    if visual_type == "flow":
        draw_card(draw, 40, 240, 370, 780, "1. User Action", ["Click / Input", "Submit / Hover", "Keyboard event"])
        draw_card(draw, 415, 240, 745, 780, "2. Handler", ["on('click', fn)", "Check data", "Call DOM update"])
        draw_card(draw, 790, 240, 1160, 780, "3. Result", ["DOM changed", "Class toggled", "User sees feedback"])
        draw.line([(370, 510), (415, 510)], fill=(40, 60, 120), width=7)
        draw.line([(745, 510), (790, 510)], fill=(40, 60, 120), width=7)
    elif visual_type == "compare":
        draw_card(draw, 60, 240, 560, 790, "jQuery", ["Short syntax", "Great in legacy", "Fast start"])
        draw_card(draw, 640, 240, 1140, 790, "Modern JavaScript", ["Native APIs", "No extra lib", "Better tree-shaking"])
        draw.text((430, 205), "VS", fill=(31, 61, 122), font=get_font(54))
    elif visual_type == "dom_tree":
        draw.rounded_rectangle([(70, 245), (1130, 790)], radius=24, fill=(248, 251, 255), outline=(110, 132, 202), width=3)
        draw.text((95, 280), "DOM Tree", fill=(32, 55, 120), font=get_font(36))
        draw.text((125, 350), "html", fill=(42, 58, 96), font=get_font(28))
        draw.text((200, 420), "body", fill=(42, 58, 96), font=get_font(28))
        draw.text((290, 500), "div#app", fill=(42, 58, 96), font=get_font(28))
        draw.text((420, 580), "ul.list", fill=(42, 58, 96), font=get_font(28))
        draw.text((560, 660), "li.item", fill=(42, 58, 96), font=get_font(28))
        draw.line([(165, 360), (228, 430)], fill=(64, 84, 146), width=5)
        draw.line([(265, 440), (325, 510)], fill=(64, 84, 146), width=5)
        draw.line([(380, 520), (458, 590)], fill=(64, 84, 146), width=5)
        draw.line([(512, 600), (596, 670)], fill=(64, 84, 146), width=5)
    else:
        draw_card(draw, 70, 250, 1130, 790, "W3Schools Focus", ["Methods", "Examples", "Practice", "Best practices"])

    img.save(path)


def ensure_slide_image(module_key: str, slide_idx: int, module_title: str, slide_title: str, visual_type: str) -> Path:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSETS_DIR / f"{module_key}_s{slide_idx + 1}.png"
    create_infographic_image(path, module_title, slide_title, visual_type)
    return path


def apply_slide_theme(slide):
    # RU: Задаем единый фон слайда для всех модулей.
    # EN: Set a unified slide background for all modules.
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["bg"]

    stripe = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(0.22),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = THEME["accent"]
    stripe.line.fill.background()


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = THEME["title"]


def add_points(slide, points):
    body = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(8.0), Inches(5.35))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(19)
        p.font.color.rgb = THEME["text"]
        p.space_after = Pt(8)


def add_footer(slide, module_title):
    footer = slide.shapes.add_textbox(Inches(0.65), Inches(6.9), Inches(12.0), Inches(0.35))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Курс: jQuery для детей 10-14 | Модуль: {module_title} | Основа: W3Schools"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(11)
    p.font.color.rgb = THEME["muted"]


def add_slide(prs, module_title, title, points, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(slide)
    add_title(slide, title)
    add_points(slide, points)
    slide.shapes.add_picture(str(image_path), Inches(9.05), Inches(1.35), width=Inches(3.9), height=Inches(4.8))
    add_footer(slide, module_title)


def create_module_presentation(file_name, module_key, module_title, slides):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # RU: В каждом модуле строго 6 слайдов по структуре урока.
    # EN: Every module has exactly 6 slides by lesson structure.
    for slide_idx, slide in enumerate(slides):
        image_path = ensure_slide_image(module_key, slide_idx, module_title, slide["title"], slide["visual"])
        add_slide(prs, module_title, slide["title"], slide["points"], image_path)

    out_path = OUTPUT_DIR / file_name
    prs.save(out_path)
    return out_path


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    modules = [
        {
            "file": "02_jquery_dom_manipulation_ru.pptx",
            "key": "dom",
            "title": "jQuery манипуляция с DOM",
            "slides": [
                {
                    "title": "Манипуляция с DOM: меняем страницу на лету",
                    "points": [
                        "DOM (Document Object Model) - это дерево элементов страницы.",
                        "По W3Schools, jQuery упрощает изменение HTML, текстов, атрибутов и CSS.",
                        "После изменения DOM пользователь сразу видит новый результат.",
                        "Мы будем работать как редакторы страницы: находим -> меняем -> проверяем.",
                        "В реальном проекте это нужно для карточек, форм, меню, уведомлений.",
                        "Главная идея: менять только то, что нужно, и не ломать структуру.",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Что изучим",
                    "points": [
                        "text() и html() - изменение текста и разметки (W3Schools jQuery HTML).",
                        "val() и attr() - работа со значениями полей и атрибутами.",
                        "addClass(), removeClass(), toggleClass() - управление состояниями UI.",
                        "css() - точечные стили, но крупные стили лучше хранить в CSS-файле.",
                        "append(), prepend(), before(), after() - вставка новых элементов.",
                        "remove(), empty() - удаление элементов и очистка содержимого.",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Ключевые объяснения",
                    "points": [
                        "text() безопаснее, если нужно вставить только обычный текст.",
                        "html() вставляет HTML-код, поэтому важно не подставлять опасные данные.",
                        "val() удобно использовать в формах: чтение и установка значения input.",
                        "attr() меняет атрибуты: src, href, title, disabled и т.д.",
                        "Классы лучше, чем css() в лоб: код чище и легче поддерживать.",
                        "Ошибка новичка: массово менять стили напрямую и путать логику с дизайном.",
                        "Лайфхак: группируй связанные изменения в одну функцию.",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Примеры и схемы (по мотивам W3Schools)",
                    "points": [
                        "$('#title').text('Привет!'); // меняем текст",
                        "$('#note').html('<b>Важно</b>'); // вставляем HTML",
                        "$('#name').val('Алия'); // задаем значение поля",
                        "$('a.main').attr('href', 'https://example.com');",
                        "$('.card').toggleClass('active'); // переключаем состояние",
                        "$('#list').append('<li>Новый пункт</li>');",
                        "Схема: Action -> jQuery method -> DOM updated.",
                    ],
                    "visual": "flow",
                },
                {
                    "title": "Мини-практика + проверка понимания",
                    "points": [
                        "Задача 1: поменяй текст заголовка и цвет кнопки через класс.",
                        "Задача 2: добавь элемент в начало списка и один в конец.",
                        "Задача 3: у формы установи значение поля val('...').",
                        "Вопрос: когда использовать text(), а когда html()?",
                        "Вопрос: почему toggleClass() удобнее, чем много css() строк?",
                        "Чек-лист: код читаемый? классы названы понятно? лишних операций нет?",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Итоги + источник",
                    "points": [
                        "Ты умеешь менять текст, HTML, атрибуты, классы и структуру DOM.",
                        "Сильный подход: минимум прямых стилей, максимум понятных классов.",
                        "Если изменить DOM аккуратно, интерфейс становится быстрым и чистым.",
                        "Подготовка к следующей теме: события и обработчики on().",
                        "Источник для углубления: W3Schools -> jQuery HTML/CSS/DOM.",
                        "Домашний мини-челлендж: сделать интерактивную карточку профиля.",
                    ],
                    "visual": "compare",
                },
            ],
        },
        {
            "file": "03_jquery_events_ru.pptx",
            "key": "events",
            "title": "jQuery события",
            "slides": [
                {
                    "title": "События: страница реагирует на действия",
                    "points": [
                        "Событие - это сигнал от пользователя или браузера (click, input, submit).",
                        "По W3Schools, jQuery события подключаются коротко и читаемо.",
                        "Правильная обработка событий делает интерфейс живым и удобным.",
                        "В реальных задачах события запускают валидацию, анимации, загрузку данных.",
                        "Логика: событие -> обработчик -> изменение DOM или состояние приложения.",
                        "Важно не плодить хаотичные обработчики в разных местах.",
                    ],
                    "visual": "flow",
                },
                {
                    "title": "Что изучим",
                    "points": [
                        "Базовые события: click, dblclick, mouseenter, mouseleave.",
                        "События форм: focus, blur, change, submit, keydown.",
                        "on() - универсальный метод для подписки на события.",
                        "off() - отключение обработчиков, чтобы не было дубликатов.",
                        "Делегирование событий: on(event, selector, handler).",
                        "Практика W3Schools: event object и preventDefault().",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Ключевые объяснения",
                    "points": [
                        "on('click', fn) - стандартный и масштабируемый подход.",
                        "submit почти всегда перехватывают для проверки формы.",
                        "preventDefault() отменяет действие по умолчанию (например, отправку формы).",
                        "event.target показывает, какой элемент вызвал событие.",
                        "Делегирование нужно, когда элементы добавляются динамически.",
                        "Ошибка новичка: ставить много обработчиков в цикле без контроля.",
                        "Лайфхак: храни обработчики рядом с логикой конкретного компонента.",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Примеры и паттерны",
                    "points": [
                        "$('#btn').on('click', () => $('.box').toggleClass('open'));",
                        "$('form').on('submit', (e) => { e.preventDefault(); ... });",
                        "$('#list').on('click', 'li', handler); // делегирование",
                        "$('input').on('keydown', (e) => console.log(e.key));",
                        "Аналогия: пульт -> кнопка -> команда устройству.",
                        "Паттерн: один вход в обработчик, ветвление по условиям внутри.",
                    ],
                    "visual": "flow",
                },
                {
                    "title": "Мини-практика и вопросы",
                    "points": [
                        "Сделай кнопку, которая переключает тему страницы.",
                        "Добавь в форму проверку: пустое поле -> показать сообщение.",
                        "Реализуй делегирование клика для динамического списка задач.",
                        "Вопрос: чем on() лучше старого click(fn) в больших проектах?",
                        "Вопрос: где полезен preventDefault(), а где нет?",
                        "Проверь: нет ли двойного срабатывания одного обработчика?",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Итоги + источник",
                    "points": [
                        "Ты научился управлять реакцией страницы на действия пользователя.",
                        "on(), event и preventDefault() - фундамент интерактивного UI.",
                        "Делегирование событий особенно важно для динамического DOM.",
                        "Грамотная архитектура событий уменьшает баги и дублирование.",
                        "Источник для углубления: W3Schools -> jQuery Events.",
                        "Следующий шаг: продвинутые операции с DOM-деревом.",
                    ],
                    "visual": "compare",
                },
            ],
        },
        {
            "file": "04_jquery_advanced_dom_ru.pptx",
            "key": "advdom",
            "title": "jQuery продвинутые операции с DOM",
            "slides": [
                {
                    "title": "Продвинутый DOM: собираем страницу как LEGO",
                    "points": [
                        "На этом уровне мы не только красим элементы, но и перестраиваем DOM.",
                        "По W3Schools, ключевые группы методов: insert, remove, traverse, clone.",
                        "Это нужно для списков задач, чатов, фильтров, карточек товаров.",
                        "Хорошая практика: менять минимальный участок дерева, а не всю страницу.",
                        "Чем чище структура DOM, тем проще поддержка и выше производительность.",
                        "Будем изучать методы как набор инструментов конструктора.",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Что изучим",
                    "points": [
                        "append(), prepend(), before(), after() - разные точки вставки.",
                        "remove(), detach(), empty() - разные сценарии удаления.",
                        "clone(true) - копирование элемента вместе с обработчиками.",
                        "parent(), parents(), children(), siblings(), next(), prev().",
                        "find(selector) - поиск внутри конкретной ветки дерева.",
                        "Как выбирать между обновлением узла и его пересозданием.",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Ключевые объяснения",
                    "points": [
                        "append() и prepend() вставляют внутрь родителя.",
                        "before() и after() вставляют рядом с текущим элементом.",
                        "detach() удаляет элемент, но сохраняет данные/обработчики для возврата.",
                        "children() смотрит только на 1 уровень, find() идет глубже.",
                        "siblings() полезен для переключателей табов и аккордеонов.",
                        "Ошибка новичка: много вложенных find(), где можно выбрать короче.",
                        "Лайфхак: кешируй контейнер var $list = $('#list').",
                    ],
                    "visual": "flow",
                },
                {
                    "title": "Примеры и стратегия обновления DOM",
                    "points": [
                        "$('#list').prepend('<li class=\"new\">Новый</li>');",
                        "$item.detach(); $('#archive').append($item);",
                        "$card.clone(true).appendTo('#cards');",
                        "$panel.find('.btn-save').prop('disabled', false);",
                        "$node.siblings().removeClass('active'); $node.addClass('active');",
                        "Стратегия: сначала выбрать корневой контейнер, потом работать внутри.",
                        "Избегай полного перерисовывания, если меняется только маленький блок.",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Мини-практика и разбор ошибок",
                    "points": [
                        "Собери TODO-список: добавление, удаление и перенос задач в архив.",
                        "Сделай вкладки: активная вкладка подсвечивается, остальные нет.",
                        "Добавь кнопку 'дублировать карточку' через clone(true).",
                        "Проверь, где children() правильнее find().",
                        "Найди лишние обращения к DOM и сократи их.",
                        "Ответь: когда detach() лучше remove()?",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Итоги + источник",
                    "points": [
                        "Ты освоил продвинутые операции: вставка, удаление, копирование, обход.",
                        "Ключ к хорошему коду: минимальные и точные изменения в DOM.",
                        "Опора на W3Schools помогает быстро сверять методы и сигнатуры.",
                        "Теперь можешь строить динамические интерфейсы с понятной структурой.",
                        "Источник: W3Schools -> jQuery Traversing / HTML / DOM Manipulation.",
                        "Финальная тема: сравнение jQuery и современного JavaScript.",
                    ],
                    "visual": "compare",
                },
            ],
        },
        {
            "file": "05_jquery_vs_modern_js_ru.pptx",
            "key": "compare",
            "title": "Сравнение jQuery и современного JavaScript",
            "slides": [
                {
                    "title": "jQuery и современный JS: кто за что отвечает",
                    "points": [
                        "Исторически jQuery решал проблемы несовместимости браузеров.",
                        "Сегодня нативный JS получил удобные API: querySelector, classList, fetch.",
                        "Но jQuery все еще встречается в legacy-проектах и быстрых прототипах.",
                        "Выбор зависит от стека команды, требований и возраста проекта.",
                        "Мы сравним не 'что лучше', а 'что уместнее в задаче'.",
                        "Подход зрелого разработчика: понимать оба инструмента.",
                    ],
                    "visual": "compare",
                },
                {
                    "title": "Что изучим",
                    "points": [
                        "Сравнение синтаксиса: селекторы, классы, события, AJAX/fetch.",
                        "Где jQuery ускоряет работу в старом коде.",
                        "Где modern JS дает лучшую производительность и размер бандла.",
                        "Как писать переходный код при миграции с jQuery на JS.",
                        "Какие практики из W3Schools полезны для обоих подходов.",
                        "Как принимать техническое решение в команде.",
                    ],
                    "visual": "cards",
                },
                {
                    "title": "Ключевые объяснения и trade-offs",
                    "points": [
                        "jQuery: короче запись и удобная цепочка методов.",
                        "Modern JS: меньше зависимостей, лучше tree-shaking, нативные API.",
                        "Для новых проектов часто достаточно чистого JavaScript.",
                        "Для старых систем резкая миграция может быть дорогой и рискованной.",
                        "Компромисс: постепенно заменять jQuery в новых модулях.",
                        "Ошибка новичка: выбирать технологию по тренду, а не по контексту.",
                        "Лайфхак: вести таблицу 'что уже мигрировано'.",
                    ],
                    "visual": "flow",
                },
                {
                    "title": "Примеры кода: jQuery vs native JS",
                    "points": [
                        "Выбор элемента: $('.item') vs document.querySelectorAll('.item').",
                        "Класс: $el.addClass('active') vs el.classList.add('active').",
                        "Событие: $btn.on('click', fn) vs btn.addEventListener('click', fn).",
                        "Запрос: $.ajax(...) vs fetch(url).then(...).",
                        "Текст: $('#x').text('Hi') vs el.textContent = 'Hi'.",
                        "Вывод: операции похожи, отличается синтаксис и зависимости.",
                    ],
                    "visual": "compare",
                },
                {
                    "title": "Мини-практика / вопросы команды",
                    "points": [
                        "Перепиши 3 jQuery-строки на modern JS.",
                        "Определи, где в старом проекте выгодно оставить jQuery.",
                        "Составь мини-план миграции на 3 этапа.",
                        "Оцени риски: поломка плагинов, время команды, тестовое покрытие.",
                        "Ответь: что важнее в вашем проекте - скорость старта или долгосрочная поддержка?",
                        "Собери аргументы 'за' и 'против' каждого подхода.",
                    ],
                    "visual": "dom_tree",
                },
                {
                    "title": "Итоги + источник",
                    "points": [
                        "jQuery не устарел полностью: он полезен в конкретных сценариях.",
                        "Modern JS стал мощным стандартом для новых проектов.",
                        "Сильный инженер сравнивает затраты, риски и поддержку.",
                        "Решение принимается по задаче, а не по личным предпочтениям.",
                        "Источник: W3Schools -> jQuery + JavaScript Tutorials.",
                        "Теперь ты умеешь аргументированно выбирать подход.",
                    ],
                    "visual": "flow",
                },
            ],
        },
    ]

    for module in modules:
        saved_path = create_module_presentation(module["file"], module["key"], module["title"], module["slides"])
        print(f"Created: {saved_path}")


if __name__ == "__main__":
    main()
